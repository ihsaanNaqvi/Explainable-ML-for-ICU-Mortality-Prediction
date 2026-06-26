"""
Generate the MS thesis defence presentation (TSU template style) — v2.
Improvements over v1:
  - Title slide now uses a real ECG/data-tensor hero image (no more fake building).
  - Slide 2 (Relevance) layout fixed — no more overlaps.
  - Novelty (slide 11) now shows three pillars: Transformer · Time-step SHAP · Clinical validation.
  - New slide 12: Reproducibility — QR codes for GitHub repo + PhysioNet dataset.
  - All real figures embedded at full quality, with captions tight to the image.
Output: d:/icu-xai/thesis/defence_presentation.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- TSU palette ----------
TSU_BLUE       = RGBColor(0x1F, 0x6A, 0xB4)
TSU_BLUE_DARK  = RGBColor(0x14, 0x4C, 0x82)
TSU_BLUE_LIGHT = RGBColor(0xD6, 0xE4, 0xF0)
TSU_GREY       = RGBColor(0x55, 0x55, 0x55)
TSU_TEXT       = RGBColor(0x22, 0x22, 0x22)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_RED     = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_GREEN   = RGBColor(0x0A, 0x87, 0x54)

FIG    = Path(r"d:/icu-xai/outputs/figures/thesis")
ASSETS = Path(r"d:/icu-xai/outputs/figures/deck_assets")
OUT    = Path(r"d:/icu-xai/thesis/defence_presentation.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = line_w
    return shp


def add_round_rect(slide, x, y, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.shadow.inherit = False
    shp.adjustments[0] = 0.10
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=TSU_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=TSU_TEXT,
                bullet_color=TSU_BLUE, line_spacing=1.20, bullet="■"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        r1 = p.add_run()
        r1.text = f"{bullet}  "
        r1.font.name = "Calibri"
        r1.font.size = Pt(size)
        r1.font.bold = True
        r1.font.color.rgb = bullet_color
        r2 = p.add_run()
        r2.text = txt
        r2.font.name = "Calibri"
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return tb


def draw_tsu_logo(slide, x, y, size=Inches(0.7)):
    """Stylised TSU shield."""
    sq = add_rect(slide, x, y, size, size, fill=RGBColor(0x0E, 0x1A, 0x2A))
    add_text(slide, x, y + Emu(int(size * 0.10)), size, Emu(int(size * 0.62)),
             "U", size=int(size / Emu(914400) * 38), bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Emu(int(size * 0.62)), size, Emu(int(size * 0.32)),
             "TOMSK · 1878",
             size=max(6, int(size / Emu(914400) * 9)), bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def page_chrome(slide, page_no, title=None, add_logo=True):
    """Left blue accent · title · top-right logo · page number · footer line."""
    add_rect(slide, 0, 0, Inches(0.18), SH, fill=TSU_BLUE)
    if title:
        add_text(slide, Inches(0.45), Inches(0.30), Inches(11.5), Inches(0.65),
                 title, size=26, bold=True, color=TSU_BLUE_DARK)
        add_rect(slide, Inches(0.45), Inches(0.92), Inches(2.6), Inches(0.06),
                 fill=TSU_BLUE)
    if add_logo:
        draw_tsu_logo(slide, Inches(12.45), Inches(0.20), size=Inches(0.7))
    add_text(slide, Inches(12.45), Inches(7.05), Inches(0.85), Inches(0.35),
             f"{page_no:02d}", size=12, bold=True, color=TSU_GREY,
             align=PP_ALIGN.RIGHT)
    add_rect(slide, Inches(0.45), Inches(7.05), Inches(11.5), Emu(9525),
             fill=TSU_BLUE_LIGHT)


def caption(slide, x, y, w, text):
    add_text(slide, x, y, w, Inches(0.35),
             text, size=10, color=TSU_GREY, italic=True)


# ============================================================
# SLIDE 1 — TITLE
# ============================================================
s = prs.slides.add_slide(BLANK)
# Full-bleed dark blue background
add_rect(s, 0, 0, SW, SH, fill=TSU_BLUE_DARK)

# Left dark column with TSU shield + label
add_rect(s, 0, 0, Inches(4.3), SH, fill=RGBColor(0x0B, 0x2F, 0x55))
draw_tsu_logo(s, Inches(0.65), Inches(0.65), size=Inches(1.4))
add_text(s, Inches(0.65), Inches(2.20), Inches(3.4), Inches(0.5),
         "NATIONAL RESEARCH", size=11, bold=True, color=RGBColor(0x9D, 0xC2, 0xE6))
add_text(s, Inches(0.65), Inches(2.55), Inches(3.4), Inches(0.5),
         "TOMSK STATE UNIVERSITY", size=14, bold=True, color=WHITE)
add_text(s, Inches(0.65), Inches(3.00), Inches(3.4), Inches(0.4),
         "Faculty of Applied Mathematics & Cybernetics",
         size=10, color=RGBColor(0xA8, 0xC4, 0xE0))

# Accent strip
add_rect(s, Inches(0.65), Inches(3.55), Inches(1.7), Inches(0.08), fill=TSU_BLUE)
add_text(s, Inches(0.65), Inches(3.65), Inches(3.4), Inches(0.4),
         "MASTER'S THESIS DEFENCE", size=12, bold=True, color=WHITE)
add_text(s, Inches(0.65), Inches(4.05), Inches(3.4), Inches(0.4),
         "01.04.02 · Applied Mathematics and Informatics",
         size=9, color=RGBColor(0xA8, 0xC4, 0xE0))

# Author block bottom
add_rect(s, Inches(0.65), Inches(5.50), Inches(0.08), Inches(1.4), fill=TSU_BLUE)
add_text(s, Inches(0.95), Inches(5.50), Inches(3.2), Inches(0.4),
         "Presented by", size=9, color=RGBColor(0xA8, 0xC4, 0xE0))
add_text(s, Inches(0.95), Inches(5.78), Inches(3.2), Inches(0.45),
         "Ihsan Naqvi", size=18, bold=True, color=WHITE)
add_text(s, Inches(0.95), Inches(6.25), Inches(3.2), Inches(0.35),
         "Supervisor:  [Supervisor name]", size=10, color=RGBColor(0xC8, 0xD9, 0xE8))
add_text(s, Inches(0.95), Inches(6.55), Inches(3.2), Inches(0.35),
         "Tomsk · 2026", size=10, color=RGBColor(0xC8, 0xD9, 0xE8))

# Hero image (right ~9 inches)
hero = ASSETS / "title_hero.png"
if hero.exists():
    s.shapes.add_picture(str(hero), Inches(4.30), 0,
                         width=Inches(9.03), height=SH)

# Title overlay box on the hero
add_rect(s, Inches(4.30), Inches(4.20), Inches(9.03), Inches(0.05),
         fill=TSU_BLUE)
add_text(s, Inches(4.55), Inches(4.40), Inches(8.6), Inches(1.4),
         "Explainable Machine Learning\nfor ICU Mortality Prediction",
         size=32, bold=True, color=WHITE, font="Calibri")
add_text(s, Inches(4.55), Inches(5.85), Inches(8.6), Inches(0.55),
         "A Time-Step SHAP Attribution Approach on PhysioNet 2012",
         size=15, color=RGBColor(0xCF, 0xE1, 0xF2), italic=True)

# Page number
add_text(s, Inches(12.55), Inches(7.10), Inches(0.7), Inches(0.30),
         "01", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 2 — RELEVANCE  (fixed layout, no overlap)
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 2, title="Relevance & Motivation")

# Single takeaway banner up top — strong, one sentence
add_round_rect(s, Inches(0.45), Inches(1.20), Inches(12.4), Inches(0.85),
               fill=TSU_BLUE_DARK)
add_text(s, Inches(0.75), Inches(1.30), Inches(12), Inches(0.65),
         "Why this matters — An ICU model must predict AND defend itself on rounds.",
         size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# Three vertical cards underneath — clean, equal heights, generous spacing
card_y = Inches(2.30)
card_h = Inches(4.55)
card_w = Inches(4.00)
xs = [Inches(0.45), Inches(4.67), Inches(8.88)]

cards = [
    {"num": "01", "head": "ICU is critical & costly",
     "big": "13–20%", "small": "of inpatient hospital costs",
     "body": "ICUs consume up to 20% of hospital budgets while caring for under 10% of admissions. "
             "A delayed escalation, missed septic shock, or wrong transfer carries irreversible cost."},
    {"num": "02", "head": "Severity scores are dated",
     "big": "1981", "small": "APACHE I  →  still in use",
     "body": "SAPS-I / APACHE collapse 24 h into a single worst-value vector and treat missingness as random. "
             "Modern monitoring makes both assumptions untenable."},
    {"num": "03", "head": "Deep models are opaque",
     "big": "0.74", "small": "probability, no reason",
     "body": "Deep nets often match clinicians on metrics but cannot defend a prediction. "
             "Trust is the bottleneck for bedside adoption, not accuracy."},
]
for i, c in enumerate(cards):
    # Card body
    add_round_rect(s, xs[i], card_y, card_w, card_h, fill=WHITE, line=TSU_BLUE)
    # Top accent bar
    add_rect(s, xs[i], card_y, card_w, Inches(0.55), fill=TSU_BLUE)
    add_text(s, xs[i] + Inches(0.22), card_y, card_w - Inches(0.4), Inches(0.55),
             c["num"], size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, xs[i] + Inches(0.85), card_y, card_w - Inches(1.0), Inches(0.55),
             c["head"], size=14, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # Big number
    add_text(s, xs[i] + Inches(0.25), card_y + Inches(0.75),
             card_w - Inches(0.5), Inches(1.0),
             c["big"], size=44, bold=True, color=TSU_BLUE_DARK)
    add_text(s, xs[i] + Inches(0.25), card_y + Inches(1.85),
             card_w - Inches(0.5), Inches(0.40),
             c["small"], size=11, color=TSU_GREY)
    # Divider
    add_rect(s, xs[i] + Inches(0.25), card_y + Inches(2.32),
             Inches(0.8), Emu(int(Pt(2))), fill=TSU_BLUE)
    # Body
    add_text(s, xs[i] + Inches(0.25), card_y + Inches(2.50),
             card_w - Inches(0.5), Inches(2.0),
             c["body"], size=11.5, color=TSU_TEXT)


# ============================================================
# SLIDE 3 — PROBLEM · HYPOTHESIS · SCOPE
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 3, title="Problem · Hypothesis · Scope")

items = [
    ("PROBLEM",
     "Build a model on PhysioNet 2012 Set-A that BOTH (a) outperforms classical severity scores "
     "and (b) explains each prediction on the (hour × variable) axis a clinician already uses."),
    ("HYPOTHESIS",
     "If we encode WHEN a measurement was taken (missingness as a parallel channel) and "
     "REDISTRIBUTE Shapley values back onto the temporal axis, the resulting attributions will "
     "agree with established clinical severity scores."),
    ("SCOPE",
     "PhysioNet 2012 Challenge · Set-A  ·  4 000 adult ICU patients  ·  first 48 h after admission  "
     "·  37 physiological + 6 descriptor variables  ·  in-hospital mortality, 13.5 % positive rate."),
]
y = Inches(1.30)
bar_h = Inches(0.60)
section_h = Inches(1.80)
for i, (label, body) in enumerate(items):
    add_round_rect(s, Inches(0.45), y, Inches(12.4), bar_h, fill=TSU_BLUE)
    add_text(s, Inches(0.75), y, Inches(12), bar_h,
             label, size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.75), y + bar_h + Inches(0.10),
             Inches(12.0), Inches(1.0),
             body, size=14, color=TSU_TEXT)
    y = y + section_h


# ============================================================
# SLIDE 4 — OBJECT · SUBJECT · AIM · TASKS
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 4, title="Object · Subject · Aim · Tasks")

# Left column
def label_block(y, label, body, height=Inches(1.4)):
    add_rect(s, Inches(0.50), y + Inches(0.05), Inches(0.15), Inches(0.50),
             fill=TSU_BLUE)
    add_text(s, Inches(0.78), y, Inches(5.8), Inches(0.45),
             label, size=15, bold=True, color=TSU_BLUE_DARK)
    add_text(s, Inches(0.78), y + Inches(0.45), Inches(5.8), height,
             body, size=12, color=TSU_TEXT)

label_block(Inches(1.30), "OBJECT",
            "In-hospital mortality prediction for adult ICU patients during their first 48 hours of admission.")
label_block(Inches(3.10), "SUBJECT",
            "Time-resolved Shapley-value attribution mapping model evidence onto the "
            "(hour × variable) axis and agreeing with established severity scores.")
label_block(Inches(5.20), "AIM",
            "Develop, evaluate, and clinically validate an explainable deep-learning predictor "
            "that matches tree-ensemble discrimination while restoring temporal interpretability.")

# Right column: 4 tasks
add_text(s, Inches(7.10), Inches(1.20), Inches(5.8), Inches(0.40),
         "Research tasks", size=18, bold=True, color=TSU_BLUE_DARK)
add_rect(s, Inches(7.10), Inches(1.65), Inches(1.0), Inches(0.05), fill=TSU_BLUE)
tasks = [
    "Build a (4000 × 48 × 37) tensor with explicit missingness encoding.",
    "Train and compare XGBoost, TCN, and Time-Aware Transformer under identical protocol.",
    "Derive a time-step SHAP procedure returning attributions on the (hour × variable) axis.",
    "Validate attributions against SOFA / SAPS-I and run a four-variant ablation.",
]
for i, t in enumerate(tasks):
    cy = Inches(1.95 + i * 1.20)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(7.10), cy, Inches(0.55), Inches(0.55))
    circ.fill.solid(); circ.fill.fore_color.rgb = TSU_BLUE
    circ.line.fill.background(); circ.shadow.inherit = False
    add_text(s, Inches(7.10), cy, Inches(0.55), Inches(0.55),
             str(i + 1), size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.85), cy - Inches(0.02), Inches(5.0), Inches(1.05),
             t, size=12, color=TSU_TEXT)


# ============================================================
# SLIDE 5 — DATA & PREPROCESSING
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 5, title="Data & Preprocessing")

# Left: stat cards
stats = [("4 000", "patients"), ("48 h", "monitoring window"),
         ("37 + 6", "vars + descriptors"), ("13.5 %", "mortality rate")]
for i, (n, u) in enumerate(stats):
    y = Inches(1.25 + i * 1.10)
    add_round_rect(s, Inches(0.55), y, Inches(2.85), Inches(0.95),
                   fill=WHITE, line=TSU_BLUE)
    add_rect(s, Inches(0.55), y, Inches(0.12), Inches(0.95), fill=TSU_BLUE)
    add_text(s, Inches(0.80), y + Inches(0.05), Inches(2.6), Inches(0.5),
             n, size=22, bold=True, color=TSU_BLUE_DARK)
    add_text(s, Inches(0.80), y + Inches(0.55), Inches(2.6), Inches(0.35),
             u, size=11, color=TSU_GREY)

# Pipeline summary
add_round_rect(s, Inches(0.55), Inches(5.85), Inches(2.85), Inches(1.05),
               fill=TSU_BLUE_LIGHT, line=TSU_BLUE)
add_text(s, Inches(0.75), Inches(5.95), Inches(2.6), Inches(0.45),
         "Final input tensor", size=12, bold=True, color=TSU_BLUE_DARK)
add_text(s, Inches(0.75), Inches(6.30), Inches(2.6), Inches(0.5),
         "(4000, 48, 74)   →  37 values + 37 masks",
         size=11, color=TSU_TEXT)

# Right: missingness figure
img = FIG / "ch3_missingness_profile.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(3.80), Inches(1.25),
                         width=Inches(9.10))
caption(s, Inches(3.80), Inches(6.55), Inches(9.10),
        "Fig. 1 — Per-variable missingness profile across 48 h. "
        "Missingness is informative and is encoded as a parallel input channel.")


# ============================================================
# SLIDE 6 — METHOD PIPELINE
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 6, title="Method — Three-Stage Pipeline")

steps = [
    ("Preprocessing",
     "Hourly median aggregation\nMissingness mask channel\nTrain/Val/Test = 2800/600/600"),
    ("Modelling",
     "XGBoost (227 feats)\nTCN (dilations 1·2·4·8)\nTime-Aware Transformer"),
    ("Time-step SHAP",
     "TreeSHAP on flat features\nRedistribute Φ over 48 hrs\nweighted by |X[h, v]|"),
    ("Clinical validation",
     "SOFA / SAPS-I concordance\nSpearman vs |calib err|\n4-variant ablation"),
]
box_w = Inches(2.85)
box_h = Inches(2.25)
y_b   = Inches(1.30)
xs    = [Inches(0.45), Inches(3.65), Inches(6.85), Inches(10.05)]

for i, (title, body) in enumerate(steps):
    add_round_rect(s, xs[i], y_b, box_w, box_h, fill=WHITE, line=TSU_BLUE)
    add_rect(s, xs[i], y_b, box_w, Inches(0.55), fill=TSU_BLUE)
    add_text(s, xs[i] + Inches(0.15), y_b,
             box_w - Inches(0.3), Inches(0.55),
             f"Step {i+1} — {title}",
             size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, xs[i] + Inches(0.18), y_b + Inches(0.65),
             box_w - Inches(0.36), Inches(1.5),
             body, size=12, color=TSU_TEXT)
    if i < 3:
        ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                xs[i] + box_w + Inches(0.02),
                                y_b + Inches(0.95),
                                Inches(0.30), Inches(0.35))
        ar.fill.solid(); ar.fill.fore_color.rgb = TSU_BLUE
        ar.line.fill.background(); ar.shadow.inherit = False

# Bottom: novelty highlight (Transformer + SHAP)
add_round_rect(s, Inches(0.45), Inches(3.90), Inches(12.4), Inches(2.95),
               fill=TSU_BLUE_LIGHT, line=TSU_BLUE)
add_text(s, Inches(0.75), Inches(4.05), Inches(12), Inches(0.5),
         "Novelty in the pipeline",
         size=18, bold=True, color=TSU_BLUE_DARK)

# Two-column novelty narrative inside the highlight box
add_text(s, Inches(0.75), Inches(4.55), Inches(5.9), Inches(0.4),
         "★  Time-Aware Transformer", size=13.5, bold=True, color=TSU_BLUE_DARK)
add_text(s, Inches(0.75), Inches(4.92), Inches(5.9), Inches(1.85),
         "Missingness-aware (48 × 74) input · learnable sinusoidal positional encoding "
         "· CLS-token head · focal loss for imbalanced ICU cohort. "
         "Designed for irregularly-sampled critical-care data, not adapted from NLP.",
         size=11.5, color=TSU_TEXT)

add_text(s, Inches(6.95), Inches(4.55), Inches(5.9), Inches(0.4),
         "★  Time-Step SHAP redistribution", size=13.5, bold=True, color=TSU_BLUE_DARK)
add_text(s, Inches(6.95), Inches(4.92), Inches(5.9), Inches(1.85),
         "TreeSHAP gives one Φ per flat feature → temporal axis is lost. "
         "We sum Φ across the 5 stats per variable, then redistribute back over 48 hours "
         "weighted by |X[h, v]|.  Conservation, locality, graceful degradation preserved.",
         size=11.5, color=TSU_TEXT)


# ============================================================
# SLIDE 7 — ARCHITECTURE
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 7, title="Architecture — Time-Aware Transformer")

img = FIG / "ch4_architecture_diagram.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.55), Inches(1.25),
                         width=Inches(8.4))

xR = Inches(9.25)
add_text(s, xR, Inches(1.25), Inches(3.8), Inches(0.4),
         "Key design choices", size=15, bold=True, color=TSU_BLUE_DARK)
add_rect(s, xR, Inches(1.65), Inches(1.0), Inches(0.05), fill=TSU_BLUE)
choices = [
    ("Input (48, 74)", "37 values + 37 missingness masks"),
    ("Time-Aware PE", "learnable sinusoidal frequencies"),
    ("Encoder", "4 layers × 4 heads (pre-norm)"),
    ("CLS head", "64 → 32 → 1 mortality logit"),
    ("Loss", "focal  γ = 2,  α = 0.861"),
    ("Params", "205 153  ·  best epoch 16"),
]
for i, (k, v) in enumerate(choices):
    y = Inches(1.95 + i * 0.78)
    add_rect(s, xR, y, Inches(0.15), Inches(0.65), fill=TSU_BLUE)
    add_text(s, xR + Inches(0.28), y, Inches(3.5), Inches(0.32),
             k, size=12, bold=True, color=TSU_BLUE_DARK)
    add_text(s, xR + Inches(0.28), y + Inches(0.30), Inches(3.5), Inches(0.40),
             v, size=11, color=TSU_TEXT)

caption(s, Inches(0.55), Inches(6.78), Inches(8.4),
        "Fig. 2 — Time-Aware Transformer.  Missingness-aware (48 × 74) input, "
        "CLS-token classification, focal loss.")


# ============================================================
# SLIDE 8 — RESULTS  (table + radar)
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 8, title="Results — Model Comparison")

# Left: comparison table
table_rows = [
    ["Model",       "AUROC", "AUPRC", "Score1", "HL-H"],
    ["SAPS-I",      "—",     "—",     "0.3097", "35.21"],
    ["XGBoost",     "0.8798","0.5384","0.5422", "37.67"],
    ["TCN",         "0.8191","0.4250","0.4253", "205.78"],
    ["Transformer", "0.8542","0.5098","0.5301", "165.17"],
]
nR, nC = len(table_rows), len(table_rows[0])
tbl = s.shapes.add_table(nR, nC,
                         Inches(0.55), Inches(1.40),
                         Inches(6.6), Inches(3.0)).table
col_w = [Inches(2.0), Inches(1.15), Inches(1.15), Inches(1.15), Inches(1.15)]
for i, w in enumerate(col_w):
    tbl.columns[i].width = w
for r in range(nR):
    tbl.rows[r].height = Inches(0.5) if r > 0 else Inches(0.55)
for r in range(nR):
    for c in range(nC):
        cell = tbl.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        run = p.add_run(); run.text = table_rows[r][c]
        run.font.name = "Calibri"; run.font.size = Pt(13)
        if r == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = TSU_BLUE_DARK
        else:
            run.font.color.rgb = TSU_TEXT
            cell.fill.solid()
            if table_rows[r][0] == "XGBoost":
                run.font.bold = True; run.font.color.rgb = TSU_BLUE_DARK
                cell.fill.fore_color.rgb = RGBColor(0xE7, 0xF4, 0xE7)
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else TSU_BLUE_LIGHT

add_text(s, Inches(0.55), Inches(4.55), Inches(6.6), Inches(0.4),
         "ML models nearly DOUBLE the SAPS-I baseline (0.31 → 0.54 on Score1).",
         size=12, bold=True, color=ACCENT_GREEN)
add_bullets(s, Inches(0.55), Inches(5.00), Inches(6.6), Inches(1.9), [
    "XGBoost — best overall: high discrimination AND calibrated.",
    "Transformer — competitive discrimination; focal loss costs calibration.",
    "Deep models become preferable at scale: longer windows, transfer.",
], size=11.5)

img = FIG / "model_radar_comparison.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(7.50), Inches(1.30),
                         width=Inches(5.4))
caption(s, Inches(7.50), Inches(6.55), Inches(5.4),
        "Fig. 3 — Radar comparison of all five discrimination metrics.")


# ============================================================
# SLIDE 9 — TIME-STEP SHAP IN ACTION
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 9, title="Novelty in Action — Time-Step SHAP")

img = FIG / "ch5_patient_shap_heatmap.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(0.55), Inches(1.25),
                         width=Inches(7.5))

xR = Inches(8.30)
add_text(s, xR, Inches(1.25), Inches(4.7), Inches(0.4),
         "What this delivers", size=15, bold=True, color=TSU_BLUE_DARK)
add_rect(s, xR, Inches(1.65), Inches(1.0), Inches(0.05), fill=TSU_BLUE)
add_bullets(s, xR, Inches(1.80), Inches(4.7), Inches(2.4), [
    "One attribution per (patient, hour, variable).",
    "Lives in the clinician's coordinate system.",
    "Exact local accuracy at the variable level.",
    "No retraining; seconds per patient.",
], size=12)

add_round_rect(s, xR, Inches(4.30), Inches(4.7), Inches(2.40),
               fill=TSU_BLUE_LIGHT, line=TSU_BLUE)
add_text(s, xR + Inches(0.15), Inches(4.40),
         Inches(4.4), Inches(0.4),
         "Top-10 attributed variables (cohort mean |Φ|)",
         size=12, bold=True, color=TSU_BLUE_DARK)
add_text(s, xR + Inches(0.15), Inches(4.75),
         Inches(4.4), Inches(1.85),
         "GCS · MechVent · Urine · BUN · Creatinine · Age · HR · MAP · Lactate · Glucose\n\n"
         "→ All 10 appear in at least one classical severity score "
         "(SAPS · APACHE · SOFA).",
         size=11.5, color=TSU_TEXT)

caption(s, Inches(0.55), Inches(6.78), Inches(7.5),
        "Fig. 4 — (hour × variable) SHAP heatmaps.  Patient A (top): non-survivor, "
        "concordance 0.80 with SAPS-I.  Patient B (bottom): survivor, concordance 0.20.")


# ============================================================
# SLIDE 10 — ABLATION + CLINICAL VALIDATION
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 10, title="Ablation & Clinical Validation")

abl = [
    ["Variant",        "Score1", "HL-H",   "ΔScore1"],
    ["full_model",     "0.5000", "257.30", "—"],
    ["no_time_pe",     "0.5301", "249.82", "+0.030"],
    ["no_missingness", "0.4337", "114.68", "−0.066"],
    ["no_focal_loss",  "0.5181", "26.31",  "+0.018"],
]
nR, nC = len(abl), len(abl[0])
tbl = s.shapes.add_table(nR, nC,
                         Inches(0.55), Inches(1.40),
                         Inches(6.4), Inches(2.7)).table
widths = [Inches(2.4), Inches(1.3), Inches(1.4), Inches(1.3)]
for i, w in enumerate(widths):
    tbl.columns[i].width = w
for r in range(nR):
    tbl.rows[r].height = Inches(0.42) if r > 0 else Inches(0.55)
for r in range(nR):
    for c in range(nC):
        cell = tbl.cell(r, c)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
        run = p.add_run(); run.text = abl[r][c]
        run.font.name = "Calibri"; run.font.size = Pt(12)
        if r == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = TSU_BLUE_DARK
        else:
            cell.fill.solid()
            if abl[r][0] == "no_missingness":
                cell.fill.fore_color.rgb = RGBColor(0xFD, 0xEA, 0xEA)
                run.font.bold = True; run.font.color.rgb = ACCENT_RED
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else TSU_BLUE_LIGHT
                run.font.color.rgb = TSU_TEXT

add_text(s, Inches(0.55), Inches(4.25), Inches(6.4), Inches(0.5),
         "Headline finding — removing missingness costs −6.6 pp on Score1.",
         size=13, bold=True, color=ACCENT_RED)
add_text(s, Inches(0.55), Inches(4.70), Inches(6.4), Inches(1.9),
         "Encoding WHEN measurements were taken is the single most informative design "
         "choice in the entire pipeline. This quantifies, on PhysioNet 2012, the long-held "
         "clinical intuition that sampling pattern is itself a signal.",
         size=12, color=TSU_TEXT)

img = FIG / "ch6_shap_sofa_concordance.png"
if img.exists():
    s.shapes.add_picture(str(img), Inches(7.20), Inches(1.30),
                         width=Inches(5.85))
add_text(s, Inches(7.20), Inches(6.10), Inches(5.85), Inches(0.4),
         "Mean SAPS-I concordance  =  0.498   (chance ≈ 0.135)",
         size=12.5, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_text(s, Inches(7.20), Inches(6.50), Inches(5.85), Inches(0.4),
         "Spearman ρ vs |calib err|  =  +0.139,   p = 6.3 × 10⁻⁴",
         size=11.5, bold=True, color=TSU_BLUE_DARK, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11 — NOVELTY (three pillars)
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 11, title="Scientific Novelty — Three Pillars")

nv = ASSETS / "novelty_visual.png"
if nv.exists():
    s.shapes.add_picture(str(nv), Inches(0.45), Inches(1.20),
                         width=Inches(12.4))

add_text(s, Inches(0.45), Inches(6.55), Inches(12.4), Inches(0.45),
         "Together: a predictor whose explanation is temporally resolved AND statistically validated.",
         size=13, bold=True, color=TSU_BLUE_DARK, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 12 — REPRODUCIBILITY (QR codes)
# ============================================================
s = prs.slides.add_slide(BLANK)
page_chrome(s, 12, title="Reproducibility & Resources")

# Two QR-code cards
def qr_card(x, png_name, title, url, accent):
    card_w = Inches(5.9)
    card_h = Inches(5.30)
    y = Inches(1.30)
    add_round_rect(s, x, y, card_w, card_h, fill=WHITE, line=accent)
    # Header
    add_rect(s, x, y, card_w, Inches(0.65), fill=accent)
    add_text(s, x + Inches(0.25), y, card_w - Inches(0.5), Inches(0.65),
             title, size=18, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    # QR
    p = ASSETS / png_name
    if p.exists():
        s.shapes.add_picture(str(p), x + Inches(1.20), y + Inches(0.95),
                             width=Inches(3.5), height=Inches(3.5))
    # URL
    add_text(s, x + Inches(0.25), y + Inches(4.55), card_w - Inches(0.5),
             Inches(0.4), url, size=10, color=TSU_TEXT, align=PP_ALIGN.CENTER)

qr_card(Inches(0.55), "qr_github.png",
        "Project repository (GitHub)",
        "github.com/ihsaanNaqvi/Explainable-ML-for-ICU-Mortality-Prediction",
        TSU_BLUE_DARK)
qr_card(Inches(6.85), "qr_physionet.png",
        "Dataset — PhysioNet 2012",
        "physionet.org/content/challenge-2012/1.0.0/",
        ACCENT_GREEN)

# Bottom bar — what to find
add_round_rect(s, Inches(0.55), Inches(6.70), Inches(12.2), Inches(0.40),
               fill=TSU_BLUE_LIGHT)
add_text(s, Inches(0.75), Inches(6.70), Inches(12), Inches(0.40),
         "Both QR codes resolve in any phone camera. Repo contains code, models, "
         "result JSONs, and the full thesis. Dataset is free under the PhysioNet license.",
         size=10.5, color=TSU_TEXT, anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13 — THANK YOU
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, fill=TSU_BLUE)
# Small shield in centre-left
draw_tsu_logo(s, Inches(1.60), Inches(2.50), size=Inches(2.6))
add_text(s, Inches(5.30), Inches(2.65), Inches(7.5), Inches(1.2),
         "THANK YOU", size=64, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(5.30), Inches(3.95), Inches(7.5), Inches(0.5),
         "Questions are welcome.", size=20, color=WHITE)
add_text(s, Inches(5.30), Inches(4.55), Inches(7.5), Inches(0.4),
         "Ihsan Naqvi  ·  TSU  ·  Tomsk, 2026", size=14,
         color=RGBColor(0xCF, 0xE1, 0xF2))
add_text(s, Inches(12.50), Inches(7.10), Inches(0.7), Inches(0.30),
         "13", size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)


# ---------- save ----------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Wrote {OUT}  -  {OUT.stat().st_size/1024:.1f} KB  -  {len(prs.slides)} slides")
